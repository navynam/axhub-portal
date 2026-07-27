# -*- coding: utf-8 -*-
"""
docs/*.md → docs/pptx/*.pptx 변환기 (python-pptx)
-----------------------------------------------------------------------------
마크다운 문서를 신한 브랜드 톤의 PowerPoint 로 변환한다.
  · 제목 슬라이드(H1 + 인용/부제)
  · H2 = 새 섹션 슬라이드, H3/H4 = 슬라이드 내 소제목
  · 불릿/번호 목록, 표, 코드블록(Mermaid 포함), 인용, 문단 지원
  · 내용이 넘치면 자동으로 다음 슬라이드로 분할("(계속)")
  · 한글은 Malgun Gothic(ea 타이포페이스) 로 설정

사용:  python docs/_md2pptx.py        (docs/*.md 전체 변환)
"""
import os, re, glob, math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 브랜드 색 ──
NAVY      = RGBColor(0x00, 0x46, 0xFF)
NAVY_DEEP = RGBColor(0x00, 0x23, 0x6E)
INK       = RGBColor(0x15, 0x22, 0x3B)
GRAY      = RGBColor(0x56, 0x5F, 0x6D)
GRAYLT    = RGBColor(0x8A, 0x94, 0xA3)
CANVAS    = RGBColor(0xF4, 0xF7, 0xFB)
LINE      = RGBColor(0xE5, 0xE9, 0xF1)
ACCENT    = RGBColor(0xC6, 0x98, 0x2E)
CODEBG    = RGBColor(0x1E, 0x24, 0x2E)
CODEFG    = RGBColor(0xE6, 0xEC, 0xF5)
HEADBG    = RGBColor(0xED, 0xF1, 0xFF)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Malgun Gothic"
MONO = "Consolas"

EMU_IN = 914400
SW, SH = Inches(13.333), Inches(7.5)
ML, MR = Inches(0.62), Inches(0.62)
BAND_H = Inches(1.12)
TOP = Inches(1.42)                 # 콘텐츠 시작 y
BOTTOM = Inches(7.02)              # 콘텐츠 하한 y
CONTENT_W = SW - ML - MR
LINE_BUDGET = 92                   # 본문 한 줄 half-width 예산(대략)


def dispw(s):
    """표시 폭(한글=2, 그 외=1)"""
    return sum(2 if ord(c) > 0x1100 else 1 for c in s)


def set_font(run, size, color=INK, bold=False, mono=False):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.color.rgb = color
    name = MONO if mono else FONT
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", name)


LINK_RE = re.compile(r"\[([^\]]+)\]\([^)]+\)")
INLINE_RE = re.compile(r"(\*\*[^*]+\*\*|`[^`]+`)")

def inline_runs(text):
    """(text, bold, mono) 런 리스트로 분해. 링크는 텍스트만, ** 볼드, ` 모노."""
    text = LINK_RE.sub(r"\1", text)
    out = []
    for part in INLINE_RE.split(text):
        if not part:
            continue
        if part.startswith("**") and part.endswith("**"):
            out.append((part[2:-2], True, False))
        elif part.startswith("`") and part.endswith("`"):
            out.append((part[1:-1], False, True))
        else:
            out.append((part, False, False))
    return out or [(text, False, False)]


# ───────────────────────── 마크다운 파서 ─────────────────────────
def parse(md):
    lines = md.split("\n")
    blocks = []
    i, n = 0, len(lines)
    while i < n:
        ln = lines[i]
        s = ln.strip()
        # 코드펜스
        if s.startswith("```"):
            lang = s[3:].strip()
            i += 1; buf = []
            while i < n and not lines[i].strip().startswith("```"):
                buf.append(lines[i]); i += 1
            i += 1
            blocks.append(("code", lang, buf)); continue
        # 빈 줄
        if not s:
            i += 1; continue
        # 수평선
        if re.match(r"^(-{3,}|\*{3,}|_{3,})$", s):
            blocks.append(("hr",)); i += 1; continue
        # 헤딩
        m = re.match(r"^(#{1,6})\s+(.*)$", s)
        if m:
            blocks.append(("h", len(m.group(1)), m.group(2).strip())); i += 1; continue
        # 표
        if s.startswith("|") and i + 1 < n and re.match(r"^\|?[\s:|-]+\|?$", lines[i+1].strip()) and "-" in lines[i+1]:
            rows = []
            def cells(row):
                r = row.strip()
                if r.startswith("|"): r = r[1:]
                if r.endswith("|"): r = r[:-1]
                return [c.strip() for c in r.split("|")]
            rows.append(cells(lines[i])); i += 2
            while i < n and lines[i].strip().startswith("|"):
                rows.append(cells(lines[i])); i += 1
            blocks.append(("table", rows)); continue
        # 불릿
        if re.match(r"^\s*[-*+]\s+", ln):
            items = []
            while i < n and re.match(r"^\s*[-*+]\s+", lines[i]):
                indent = len(lines[i]) - len(lines[i].lstrip())
                txt = re.sub(r"^\s*[-*+]\s+", "", lines[i]).strip()
                items.append((min(indent // 2, 2), txt)); i += 1
            blocks.append(("bullets", items)); continue
        # 번호목록
        if re.match(r"^\s*\d+\.\s+", ln):
            items = []
            while i < n and re.match(r"^\s*\d+\.\s+", lines[i]):
                txt = re.sub(r"^\s*\d+\.\s+", "", lines[i]).strip()
                items.append(txt); i += 1
            blocks.append(("numbered", items)); continue
        # 인용
        if s.startswith(">"):
            buf = []
            while i < n and lines[i].strip().startswith(">"):
                buf.append(re.sub(r"^\s*>\s?", "", lines[i])); i += 1
            blocks.append(("quote", " ".join(x.strip() for x in buf if x.strip()))); continue
        # 문단
        buf = []
        while i < n and lines[i].strip() and not re.match(r"^(#{1,6}\s|```|\||\s*[-*+]\s|\s*\d+\.\s|>)", lines[i]) \
                and not re.match(r"^(-{3,}|\*{3,})$", lines[i].strip()):
            buf.append(lines[i].strip()); i += 1
        blocks.append(("para", " ".join(buf)))
    return blocks


# ───────────────────────── 덱 빌더 ─────────────────────────
class Deck:
    def __init__(self, title, subtitle, fname):
        self.prs = Presentation()
        self.prs.slide_width = SW; self.prs.slide_height = SH
        self.blank = self.prs.slide_layouts[6]
        self.title = title
        self.fname = fname
        self.section = title
        self.slide = None
        self.y = 0
        self.pageno = 0
        self._title_slide(title, subtitle)

    # 배경
    def _bg(self, slide, color=WHITE):
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
        r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background()
        r.shadow.inherit = False
        slide.shapes._spTree.remove(r._element); slide.shapes._spTree.insert(2, r._element)
        return r

    def _rect(self, slide, l, t, w, h, color, lncolor=None):
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        r.fill.solid(); r.fill.fore_color.rgb = color
        if lncolor is None:
            r.line.fill.background()
        else:
            r.line.color.rgb = lncolor; r.line.width = Pt(0.75)
        r.shadow.inherit = False
        return r

    def _tb(self, slide, l, t, w, h):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
        tf.margin_top = 0; tf.margin_bottom = 0
        return tf

    def _para(self, tf, runs, size, color=INK, bold=False, align=PP_ALIGN.LEFT,
              first=False, space_before=0.0, indent=0.0, bullet=None, line=1.0):
        p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2); p.space_before = Pt(space_before)
        if line != 1.0:
            p.line_spacing = line
        if indent:
            p.paragraph_format.left_indent = Inches(indent)
        if bullet:
            r = p.add_run(); r.text = bullet
            set_font(r, size, color=NAVY if bullet.strip() == "•" else GRAYLT, bold=False)
        for (t, b, mono) in runs:
            if t == "":
                continue
            r = p.add_run(); r.text = t
            set_font(r, size, color=color, bold=(bold or b), mono=mono)
        return p

    def _title_slide(self, title, subtitle):
        s = self.prs.slides.add_slide(self.blank)
        self._bg(s, WHITE)
        self._rect(s, 0, 0, Inches(0.28), SH, NAVY)              # 좌측 액센트 바
        self._rect(s, Inches(0.9), Inches(2.55), Inches(1.4), Inches(0.09), ACCENT)
        tf = self._tb(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(2.2))
        self._para(tf, inline_runs(title), 40, color=INK, bold=True, first=True, line=1.1)
        if subtitle:
            tf2 = self._tb(s, Inches(0.92), Inches(4.15), Inches(11.3), Inches(1.6))
            for i, line in enumerate(subtitle):
                self._para(tf2, inline_runs(line), 15, color=GRAY, first=(i == 0), line=1.3)
        # 푸터
        tf3 = self._tb(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.4))
        self._para(tf3, [("AX-HUB · 신한라이프   |   %s" % self.fname, False, False)], 11, color=GRAYLT, first=True)

    def _new_slide(self, cont=False):
        if self.slide is not None:
            self._footer()
        s = self.prs.slides.add_slide(self.blank)
        self._bg(s, WHITE)
        self._rect(s, 0, 0, SW, BAND_H, NAVY)
        self._rect(s, 0, BAND_H, SW, Inches(0.05), ACCENT)
        head = self.section + ("  · (계속)" if cont else "")
        tf = self._tb(s, ML, Inches(0.24), CONTENT_W, Inches(0.7))
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._para(tf, inline_runs(head), 21, color=WHITE, bold=True, first=True, line=1.0)
        self.slide = s
        self.y = int(TOP)
        self.pageno += 1

    def _footer(self):
        tf = self._tb(self.slide, ML, Inches(7.08), CONTENT_W, Inches(0.32))
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = self.title; set_font(r, 9, color=GRAYLT)
        tf2 = self._tb(self.slide, SW - Inches(1.4), Inches(7.08), Inches(0.9), Inches(0.32))
        p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
        r2 = p2.add_run(); r2.text = str(self.pageno); set_font(r2, 9, color=GRAYLT)

    def ensure(self, need):
        if self.slide is None or self.y + need > int(BOTTOM):
            self._new_slide(cont=(self.slide is not None))

    # ── 블록 렌더 ──
    def heading(self, level, text):
        if level <= 2:
            self.section = text
            self._new_slide(cont=False)
        else:
            need = Inches(0.62)
            self.ensure(need)
            tf = self._tb(self.slide, ML, self.y, CONTENT_W, Inches(0.5))
            pref = "▸ " if level == 3 else "· "
            self._para(tf, [(pref, False, False)] + inline_runs(text), 17 if level == 3 else 15,
                       color=NAVY if level == 3 else INK, bold=True, first=True)
            self.y += int(Inches(0.5))

    def para(self, text):
        if not text.strip():
            return
        lines = math.ceil(dispw(text) / LINE_BUDGET)
        h = Inches(0.02) + Inches(0.30) * lines
        self.ensure(h)
        tf = self._tb(self.slide, ML, self.y, CONTENT_W, h)
        self._para(tf, inline_runs(text), 14.5, color=INK, first=True, line=1.15)
        self.y += int(h) + int(Inches(0.06))

    def quote(self, text):
        lines = math.ceil(dispw(text) / (LINE_BUDGET - 8))
        h = Inches(0.16) + Inches(0.28) * lines
        self.ensure(h + Inches(0.1))
        bar = self._rect(self.slide, ML, self.y, Inches(0.06), h, ACCENT)
        box = self._rect(self.slide, ML + Inches(0.06), self.y, CONTENT_W - Inches(0.06), h, HEADBG)
        tf = self._tb(self.slide, ML + Inches(0.24), self.y + int(Inches(0.06)), CONTENT_W - Inches(0.4), h - Inches(0.12))
        self._para(tf, inline_runs(text), 13, color=GRAY, first=True, line=1.15)
        self.y += int(h) + int(Inches(0.1))

    def bullets(self, items):
        for lvl, txt in items:
            lines = math.ceil(dispw(txt) / (LINE_BUDGET - 6 - lvl * 6))
            h = Inches(0.04) + Inches(0.30) * lines
            self.ensure(h)
            tf = self._tb(self.slide, ML + Inches(0.1 + lvl * 0.32), self.y, CONTENT_W - Inches(0.1 + lvl * 0.32), h)
            bullet = "•  " if lvl == 0 else ("–  " if lvl == 1 else "·  ")
            self._para(tf, inline_runs(txt), 14, color=INK, first=True, bullet=bullet, line=1.1)
            self.y += int(h)

    def numbered(self, items):
        for idx, txt in enumerate(items, 1):
            lines = math.ceil(dispw(txt) / (LINE_BUDGET - 8))
            h = Inches(0.04) + Inches(0.30) * lines
            self.ensure(h)
            tf = self._tb(self.slide, ML + Inches(0.1), self.y, CONTENT_W - Inches(0.1), h)
            self._para(tf, [("%d.  " % idx, True, False)] + inline_runs(txt), 14, color=INK, first=True, line=1.1)
            self.y += int(h)

    def code(self, lang, lines):
        label = ("Mermaid 다이어그램 (소스)" if lang == "mermaid" else (lang or "code"))
        buf = [ln.rstrip() for ln in lines]
        idx = 0
        first_chunk = True
        while idx < len(buf):
            self.ensure(Inches(0.9))
            avail = int(BOTTOM) - self.y - int(Inches(0.4))
            per = int(Inches(0.205))
            cap = max(3, avail // per)
            chunk = buf[idx: idx + cap]; idx += cap
            # 라벨
            tfl = self._tb(self.slide, ML, self.y, CONTENT_W, Inches(0.26))
            self._para(tfl, [((label if first_chunk else label + " (계속)"), False, True)], 9.5, color=GRAYLT, first=True)
            self.y += int(Inches(0.26))
            h = per * len(chunk) + int(Inches(0.18))
            box = self._rect(self.slide, ML, self.y, CONTENT_W, h, CODEBG)
            tf = self._tb(self.slide, ML + Inches(0.16), self.y + int(Inches(0.09)), CONTENT_W - Inches(0.3), h - Inches(0.18))
            for j, cl in enumerate(chunk):
                self._para(tf, [(cl if cl else " ", False, True)], 9.5, color=CODEFG, first=(j == 0), line=1.0)
            self.y += int(h) + int(Inches(0.1))
            first_chunk = False

    def table(self, rows):
        if not rows:
            return
        ncol = max(len(r) for r in rows)
        rows = [r + [""] * (ncol - len(r)) for r in rows]
        header, body = rows[0], rows[1:]
        # 열 폭(내용 최대폭 비례)
        widths = []
        for c in range(ncol):
            mx = max(dispw(rows[r][c]) for r in range(len(rows)))
            widths.append(max(6, min(mx, 60)))
        tot = sum(widths)
        col_emu = [max(int(Inches(0.9)), int(CONTENT_W * w / tot)) for w in widths]
        # 폭 합 보정
        scale = int(CONTENT_W) / sum(col_emu)
        col_emu = [int(x * scale) for x in col_emu]
        col_hw = [max(6, int(LINE_BUDGET * e / int(CONTENT_W))) for e in col_emu]

        def row_h(cells, base):
            ln = 1
            for c in range(ncol):
                ln = max(ln, math.ceil(max(1, dispw(cells[c])) / max(4, col_hw[c])))
            return int(Inches(0.10)) + int(Inches(base)) * ln

        r = 0
        first_chunk = True
        while r < len(body) or first_chunk:
            hh = row_h(header, 0.30)
            self.ensure(hh + int(Inches(0.5)))
            # 이 슬라이드에 담을 행 수 계산
            avail = int(BOTTOM) - self.y
            used = hh
            take = []
            while r < len(body):
                rh = row_h(body[r], 0.28)
                if used + rh > avail and take:
                    break
                take.append(body[r]); used += rh; r += 1
                if used > avail:
                    break
            data = [header] + take
            total_h = hh + sum(row_h(x, 0.28) for x in take)
            gtbl = self.slide.shapes.add_table(len(data), ncol, ML, self.y, int(CONTENT_W), total_h)
            tbl = gtbl.table
            tbl.first_row = True; tbl.horz_banding = True
            for c in range(ncol):
                tbl.columns[c].width = col_emu[c]
            for ri, rowcells in enumerate(data):
                tbl.rows[ri].height = row_h(rowcells, 0.30 if ri == 0 else 0.28)
                for ci in range(ncol):
                    cell = tbl.cell(ri, ci)
                    cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
                    cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    if ri == 0:
                        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
                    else:
                        cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else CANVAS
                    tf = cell.text_frame; tf.word_wrap = True
                    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
                    for (t, b, mono) in inline_runs(rowcells[ci]):
                        rr = p.add_run(); rr.text = t
                        set_font(rr, 10.5, color=(WHITE if ri == 0 else INK), bold=(ri == 0 or b), mono=mono)
            self.y += total_h + int(Inches(0.12))
            first_chunk = False
            if r >= len(body):
                break

    def finish(self, path):
        if self.slide is not None:
            self._footer()
        self.prs.save(path)


# ───────────────────────── 실행 ─────────────────────────
def convert(md_path, out_path):
    with open(md_path, encoding="utf-8") as f:
        md = f.read()
    blocks = parse(md)
    # 제목/부제 추출
    title = os.path.splitext(os.path.basename(md_path))[0]
    subtitle = []
    start = 0
    if blocks and blocks[0][0] == "h" and blocks[0][1] == 1:
        title = blocks[0][2]; start = 1
        j = start
        while j < len(blocks) and blocks[j][0] in ("quote", "para"):
            subtitle.append(blocks[j][1]); start = j + 1; j += 1
            if len(subtitle) >= 3:
                break
    deck = Deck(title, subtitle, os.path.basename(md_path))
    for blk in blocks[start:]:
        k = blk[0]
        if k == "h":
            deck.heading(blk[1], blk[2])
        elif k == "para":
            deck.para(blk[1])
        elif k == "quote":
            deck.quote(blk[1])
        elif k == "bullets":
            deck.bullets(blk[1])
        elif k == "numbered":
            deck.numbered(blk[1])
        elif k == "code":
            deck.code(blk[1], blk[2])
        elif k == "table":
            deck.table(blk[1])
        elif k == "hr":
            pass
    deck.finish(out_path)
    return len(deck.prs.slides._sldIdLst)


def main():
    here = os.path.dirname(os.path.abspath(__file__))
    outdir = os.path.join(here, "pptx")
    os.makedirs(outdir, exist_ok=True)
    total = 0
    for md in sorted(glob.glob(os.path.join(here, "*.md"))):
        name = os.path.splitext(os.path.basename(md))[0]
        out = os.path.join(outdir, name + ".pptx")
        n = convert(md, out)
        total += 1
        print("  ok  %-22s -> pptx/%s.pptx  (%d slides)" % (os.path.basename(md), name, n))
    print("DONE %d files" % total)


if __name__ == "__main__":
    main()
